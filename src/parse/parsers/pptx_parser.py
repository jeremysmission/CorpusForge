"""
PPTX parser -- reads modern Microsoft PowerPoint (.pptx) presentations.

Plain English: opens a slide deck and pulls the text from every shape
on every slide. Each slide is tagged ``[SLIDE N]`` so operators can
trace chunks back to the slide they came from. Uses the ``python-pptx``
library; if the file is broken or the library is unavailable the parser
logs the issue and returns empty text.
"""

from __future__ import annotations

import logging
from pathlib import Path

from src.parse.parsers.txt_parser import ParsedDocument

logger = logging.getLogger(__name__)


class PptxParser:
    """Extract slide text from modern PowerPoint .pptx decks."""

    def parse(self, file_path: Path) -> ParsedDocument:
        """Open a .pptx deck and return all slide text in reading order."""
        path = Path(file_path)
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [f"[SLIDE {i}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                if len(parts) > 1:
                    slides.append("\n".join(parts))
            text = "\n\n".join(slides)
        except Exception as e:
            logger.error("PPTX parse failed for %s: %s", path.name, e)
            text = ""

        return ParsedDocument(
            source_path=str(path),
            text=text,
            parse_quality=0.9 if text.strip() else 0.0,
            file_ext=".pptx",
            file_size=path.stat().st_size,
        )
